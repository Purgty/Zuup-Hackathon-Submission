import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "ZUUP: Autonomous Public Transport AI"
    subtitle.text = "Solving the Rigid Transport Problem with Dynamic Rerouting\nHackathon Submission"
    
    # Folder containing markdown files
    md_folder = "project_exp"
    files = [
        "01_Problem_and_Solution.md",
        "02_Core_Algorithms.md",
        "03_Edge_Cases.md",
        "04_Architecture_and_Tech_Stack.md",
        "05_Business_Model.md",
        "06_Simulation_Scenarios_Deep_Dive.md"
    ]
    
    for filename in files:
        filepath = os.path.join(md_folder, filename)
        if not os.path.exists(filepath):
            continue
            
        with open(filepath, 'r', encoding='utf-8') as f:
            lines = f.readlines()
            
        # Parse markdown into slide blocks
        slide_title = ""
        content_lines = []
        
        # A simple parser: headers (## or #) become slide titles or bullet points
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith("# ") or line.startswith("## ") or line.startswith("### "):
                # If we already have a title and content, create a slide
                if slide_title and content_lines:
                    add_slide(prs, slide_title, content_lines)
                    content_lines = []
                
                # New title
                slide_title = line.replace("#", "").strip()
            else:
                # Clean markdown formatting (simple)
                line = line.replace("**", "")
                if line.startswith("- "):
                    content_lines.append(line)
                else:
                    content_lines.append("• " + line)
        
        # Add the last block
        if slide_title and content_lines:
            add_slide(prs, slide_title, content_lines)

    # Save
    prs.save("ZUUP_Pitch_Deck.pptx")
    print("Successfully generated ZUUP_Pitch_Deck.pptx")

def add_slide(prs, title_text, content_lines):
    # Bullet slide layout
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title_text
    
    tf = body_shape.text_frame
    tf.clear()
    
    for i, line in enumerate(content_lines):
        p = tf.add_paragraph()
        # Set level based on indentation
        if line.startswith("    -") or line.startswith("  -"):
            p.level = 1
            line = line.strip().replace("- ", "")
        elif line.startswith("- "):
            p.level = 0
            line = line[2:]
        elif line.startswith("• "):
            p.level = 0
            line = line[2:]
            
        p.text = line
        p.font.size = Pt(16)
        # Add some spacing
        p.space_after = Pt(10)

if __name__ == "__main__":
    create_presentation()
